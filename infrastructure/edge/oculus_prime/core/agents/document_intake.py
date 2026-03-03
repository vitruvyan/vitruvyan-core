"""
Vitruvyan INTAKE — Document Intake Agent

Media Scope: PDF, DOCX, MD, TXT, JSON, XML
Constraints: NO semantic inference, NO relevance judgment

Compliance: ACCORDO-FONDATIVO-INTAKE-V1.1

Key Principles:
- Extract text literally (descriptive, not interpretative)
- Preserve raw reference + hash
- Emit Evidence Pack + event
- NO domain-specific logic
"""

import uuid
import hashlib
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, Any, Optional, List
import logging

from infrastructure.edge.oculus_prime.core.guardrails import IntakeGuardrails

# Document processing libraries (conditional imports)
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False


logger = logging.getLogger(__name__)


class DocumentIntakeAgent:
    """
    Pre-epistemic document acquisition agent
    
    Responsibilities:
    - Extract text from documents (PDF, DOCX, MD, TXT, JSON, XML)
    - Generate normalized_text (descriptive-literal only)
    - Create immutable Evidence Pack
    - Emit oculus_prime.evidence.created event (legacy alias: intake.evidence.created)
    
    DOES NOT:
    - Interpret content semantically
    - Evaluate relevance or importance
    - Apply domain-specific rules
    - Call Codex directly
    """
    
    SUPPORTED_FORMATS = ['.pdf', '.docx', '.xlsx', '.pptx', '.md', '.txt', '.json', '.xml']
    AGENT_ID = "doc-intake-v1"
    AGENT_VERSION = "1.0.0"
    
    def __init__(self, event_emitter, postgres_agent=None):
        """
        Args:
            event_emitter: IntakeEventEmitter instance
            postgres_agent: PostgresAgent for Evidence Pack persistence
        """
        self.event_emitter = event_emitter
        self.postgres_agent = postgres_agent
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check availability of document processing libraries"""
        missing = []
        if not PDFPLUMBER_AVAILABLE:
            missing.append("pdfplumber")
        if not DOCX_AVAILABLE:
            missing.append("python-docx")
        if not OPENPYXL_AVAILABLE:
            missing.append("openpyxl")
        if not PPTX_AVAILABLE:
            missing.append("python-pptx")
        if not MARKDOWN_AVAILABLE:
            missing.append("markdown")
        
        if missing:
            logger.warning(f"Missing optional dependencies: {', '.join(missing)}")
    
    def ingest_document(
        self,
        source_path: str,
        chunking_strategy: str = "none",
        chunk_size: int = 4000,
        sampling_policy_ref: Optional[str] = None,
        correlation_id: Optional[str] = None,
        tenant_id: Optional[str] = None,
        project_name: Optional[str] = None,
    ) -> List[str]:
        """
        Ingest document and create Evidence Pack(s)
        
        Args:
            source_path: Path to document file
            chunking_strategy: 'none', 'page', 'size' (character count)
            chunk_size: Max characters per chunk (if strategy='size')
            sampling_policy_ref: Reference to external Sampling Policy
            correlation_id: Distributed tracing ID
            tenant_id: Tenant isolation key (propagated to evidence pack + events)
        
        Returns:
            List of evidence_ids created
        
        Raises:
            ValueError: If file format unsupported or unreadable
        """
        source_path = Path(source_path)
        
        # Validate file exists and format supported
        if not source_path.exists():
            raise ValueError(f"Source file not found: {source_path}")
        
        if source_path.suffix.lower() not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {source_path.suffix}. Supported: {self.SUPPORTED_FORMATS}")
        
        # Compute source hash
        source_hash = self._compute_file_hash(source_path)
        byte_size = source_path.stat().st_size
        mime_type = self._get_mime_type(source_path)
        
        # Extract text
        try:
            extracted_text = self._extract_text(source_path)
        except Exception as e:
            logger.error(f"Text extraction failed for {source_path}: {e}")
            raise ValueError(f"Text extraction failed: {e}")
        
        # Detect language (simple heuristic)
        language_detected = self._detect_language_simple(extracted_text)
        
        # Chunking
        chunks = self._chunk_text(extracted_text, chunking_strategy, chunk_size)
        
        # Create Evidence Packs
        evidence_ids = []
        for chunk_idx, chunk_text in enumerate(chunks):
            evidence_id = self._generate_evidence_id()
            chunk_id = f"CHK-{chunk_idx}"
            
            # Build Evidence Pack
            evidence_pack = {
                "evidence_id": evidence_id,
                "chunk_id": chunk_id,
                "schema_version": "1.0.0",
                "created_utc": datetime.now(timezone.utc).isoformat(),
                "source_ref": {
                    "source_type": "document",
                    "source_uri": str(source_path.absolute()),
                    "source_hash": source_hash,
                    "mime_type": mime_type,
                    "byte_size": byte_size
                },
                "normalized_text": chunk_text,
                "technical_metadata": {
                    "extraction_method": self._get_extraction_method(source_path),
                    "extraction_version": self.AGENT_VERSION,
                    "language_detected": language_detected,
                    "confidence_score": 1.0,  # Text extraction is deterministic
                    "chunk_position": {
                        "start_offset": chunk_idx * chunk_size if chunking_strategy == "size" else 0,
                        "end_offset": (chunk_idx + 1) * chunk_size if chunking_strategy == "size" else len(chunk_text),
                        "total_chunks": len(chunks)
                    }
                },
                "integrity": {
                    "evidence_hash": self._compute_evidence_hash(evidence_id, chunk_id, chunk_text, source_hash),
                    "immutable": True
                },
                "sampling_policy_ref": sampling_policy_ref,
                "tags": []
            }
            
            # Tenant isolation: include tenant_id if provided
            if tenant_id:
                evidence_pack["tenant_id"] = tenant_id
            
            # ✅ HARDENING: Validate no semantic interpretation
            try:
                IntakeGuardrails.validate_no_semantics(chunk_text, "document")
                IntakeGuardrails.validate_source_hash_required(evidence_pack["source_ref"])
            except Exception as e:
                logger.error(f"Guardrail violation: {e}")
                raise
            
            # Persist Evidence Pack (append-only)
            evidence_pack_ref = self._persist_evidence_pack(evidence_pack)
            
            # Emit event
            self.event_emitter.emit_evidence_created(
                evidence_id=evidence_id,
                chunk_id=chunk_id,
                source_type="document",
                source_uri=str(source_path.absolute()),
                evidence_pack_ref=evidence_pack_ref,
                source_hash=source_hash,
                intake_agent_id=self.AGENT_ID,
                intake_agent_version=self.AGENT_VERSION,
                byte_size=byte_size,
                language_detected=language_detected,
                sampling_policy_ref=sampling_policy_ref,
                correlation_id=correlation_id,
                tenant_id=tenant_id,
                project_name=project_name,
            )
            
            evidence_ids.append(evidence_id)
        
        return evidence_ids
    
    def _extract_text(self, source_path: Path) -> str:
        """
        Extract text from document (format-specific)
        
        Returns descriptive-literal text only. NO interpretation.
        """
        suffix = source_path.suffix.lower()
        
        if suffix == '.pdf':
            return self._extract_pdf(source_path)
        elif suffix == '.docx':
            return self._extract_docx(source_path)
        elif suffix == '.xlsx':
            return self._extract_xlsx(source_path)
        elif suffix == '.pptx':
            return self._extract_pptx(source_path)
        elif suffix == '.md':
            return self._extract_markdown(source_path)
        elif suffix == '.txt':
            return self._extract_txt(source_path)
        elif suffix == '.json':
            return self._extract_json(source_path)
        elif suffix == '.xml':
            return self._extract_xml(source_path)
        else:
            raise ValueError(f"No extraction method for: {suffix}")
    
    def _extract_pdf(self, source_path: Path) -> str:
        """Extract text from PDF (literal only)"""
        if not PDFPLUMBER_AVAILABLE:
            raise ValueError("pdfplumber not installed. Cannot process PDF.")
        
        text_parts = []
        with pdfplumber.open(source_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        return "\n\n".join(text_parts)
    
    def _extract_docx(self, source_path: Path) -> str:
        """Extract text from DOCX (literal only)"""
        if not DOCX_AVAILABLE:
            raise ValueError("python-docx not installed. Cannot process DOCX.")
        
        doc = DocxDocument(source_path)
        return "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    
    def _extract_xlsx(self, source_path: Path) -> str:
        """Extract text from XLSX — all sheets, all cell values (literal only)"""
        if not OPENPYXL_AVAILABLE:
            raise ValueError("openpyxl not installed. Cannot process XLSX.")
        
        wb = openpyxl.load_workbook(source_path, data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).rstrip()
                if line:
                    parts.append(line)
        wb.close()
        return "\n".join(parts)
    
    def _extract_pptx(self, source_path: Path) -> str:
        """Extract text from PPTX — all slides, all text frames (literal only)"""
        if not PPTX_AVAILABLE:
            raise ValueError("python-pptx not installed. Cannot process PPTX.")
        
        prs = PptxPresentation(source_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts)
    
    def _extract_markdown(self, source_path: Path) -> str:
        """Extract text from Markdown (preserve structure)"""
        with open(source_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_txt(self, source_path: Path) -> str:
        """Extract text from TXT (as-is)"""
        with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    
    def _extract_json(self, source_path: Path) -> str:
        """Extract text from JSON (structured representation)"""
        with open(source_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Return pretty-printed JSON as normalized text
        return json.dumps(data, indent=2)
    
    def _extract_xml(self, source_path: Path) -> str:
        """Extract text from XML (as-is)"""
        with open(source_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _chunk_text(self, text: str, strategy: str, chunk_size: int, overlap: int = 200) -> List[str]:
        """
        Chunk text according to strategy
        
        Strategies:
        - 'none': No chunking (single chunk)
        - 'size': Sentence-aware split by character count with overlap
        - 'page': Not applicable for text (fallback to 'none')
        
        Args:
            text: Full text to chunk
            strategy: Chunking strategy
            chunk_size: Target max characters per chunk
            overlap: Character overlap between consecutive chunks (default 200)
        """
        if strategy == "none" or strategy == "page":
            return [text]
        elif strategy == "size":
            if not text:
                return [""]
            return self._sentence_aware_chunk(text, chunk_size, overlap)
        else:
            raise ValueError(f"Unknown chunking strategy: {strategy}")

    @staticmethod
    def _sentence_aware_chunk(text: str, chunk_size: int, overlap: int) -> List[str]:
        """
        Split text into chunks at sentence boundaries with overlap.
        
        Avoids cutting sentences in half. Falls back to character split
        only for very long sentences that exceed chunk_size.
        """
        import re
        # Split on sentence-ending punctuation followed by whitespace
        sentence_pattern = re.compile(r'(?<=[.!?;\n])\s+')
        sentences = sentence_pattern.split(text)
        
        if not sentences:
            return [text] if text else [""]
        
        chunks: List[str] = []
        current_chunk = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            # If adding this sentence exceeds chunk_size and we already have content
            if current_chunk and len(current_chunk) + len(sentence) + 1 > chunk_size:
                chunks.append(current_chunk.strip())
                # Start next chunk with overlap from the end of current chunk
                if overlap > 0 and len(current_chunk) > overlap:
                    # Take last N chars, but try to start at a sentence boundary
                    overlap_text = current_chunk[-overlap:]
                    # Find first sentence start in overlap region
                    boundary = re.search(r'(?<=[.!?;])\s+', overlap_text)
                    if boundary:
                        overlap_text = overlap_text[boundary.end():]
                    current_chunk = overlap_text + " " + sentence
                else:
                    current_chunk = sentence
            else:
                current_chunk = (current_chunk + " " + sentence).strip() if current_chunk else sentence
            
            # Safety: if a single sentence exceeds chunk_size, force-split it
            while len(current_chunk) > chunk_size * 1.5:
                split_point = chunk_size
                # Try to split at a space
                space_pos = current_chunk.rfind(' ', 0, chunk_size)
                if space_pos > chunk_size * 0.5:
                    split_point = space_pos
                chunks.append(current_chunk[:split_point].strip())
                current_chunk = current_chunk[split_point:].strip()
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks if chunks else [text]
    
    def _detect_language_simple(self, text: str) -> str:
        """
        Word-frequency language detection (no external dependencies).
        
        Detects: Italian, English, French, Spanish, German.
        Uses high-frequency function-word signatures per language.
        Returns ISO 639-1 code (e.g., 'en', 'it', 'fr', 'es', 'de').
        Fallback: 'unknown'
        """
        if not text or not text.strip():
            return "unknown"
        
        # Lowercase sample (first 5000 chars for speed)
        sample = text[:5000].lower()
        # Extract words (letters only, 2+ chars)
        import re
        words = re.findall(r'\b[a-zà-öø-ÿ]{2,}\b', sample)
        if len(words) < 10:
            return "unknown"
        
        word_set = set(words)
        
        # High-frequency function words per language (articles, prepositions, conjunctions)
        lang_markers = {
            "it": {"il", "lo", "la", "le", "gli", "del", "della", "delle", "dei",
                   "nel", "nella", "che", "per", "con", "sono", "una", "uno",
                   "questo", "questa", "anche", "essere", "può", "alla", "alle",
                   "dalla", "delle", "tra", "fra", "come", "dove", "ogni",
                   "quali", "ciascun", "nonché", "inoltre", "articolo", "comma",
                   "previsto", "misure", "rischio", "sicurezza", "normativa"},
            "en": {"the", "and", "for", "with", "that", "this", "from", "have",
                   "has", "are", "was", "were", "been", "will", "would", "should",
                   "which", "their", "they", "them", "than", "into", "each",
                   "between", "through", "during", "before", "after", "shall"},
            "fr": {"les", "des", "une", "dans", "pour", "par", "sur", "avec",
                   "qui", "que", "sont", "aux", "cette", "ses", "ont", "été",
                   "entre", "depuis", "leur", "mais", "aussi", "peut"},
            "es": {"los", "las", "del", "una", "para", "por", "con", "que",
                   "son", "como", "más", "esta", "entre", "sus", "pero",
                   "desde", "cada", "han", "puede", "también", "sobre"},
            "de": {"der", "die", "das", "den", "dem", "ein", "eine", "und",
                   "für", "mit", "von", "auf", "ist", "sind", "wird", "des",
                   "auch", "nach", "sich", "bei", "oder", "über", "als"}
        }
        
        scores = {}
        for lang, markers in lang_markers.items():
            hits = word_set & markers
            # Weight by how many marker-word occurrences appear in the text
            count = sum(1 for w in words if w in markers)
            scores[lang] = (len(hits), count)
        
        # Pick language with most distinct marker hits, break ties by occurrence count
        best = max(scores.items(), key=lambda x: (x[1][0], x[1][1]))
        
        # Require at least 3 distinct marker words to be confident
        if best[1][0] < 3:
            return "unknown"
        
        return best[0]
    
    def _compute_file_hash(self, source_path: Path) -> str:
        """Compute SHA-256 hash of file"""
        sha256 = hashlib.sha256()
        with open(source_path, 'rb') as f:
            while chunk := f.read(8192):
                sha256.update(chunk)
        return f"sha256:{sha256.hexdigest()}"
    
    def _compute_evidence_hash(self, evidence_id: str, chunk_id: str, text: str, source_hash: str) -> str:
        """Compute evidence integrity hash"""
        composite = f"{evidence_id}{chunk_id}{text}{source_hash}"
        return f"sha256:{hashlib.sha256(composite.encode('utf-8')).hexdigest()}"
    
    def _generate_evidence_id(self) -> str:
        """Generate unique Evidence ID"""
        uid = uuid.uuid4()
        return f"EVD-{uid.hex.upper()[:8]}-{uid.hex.upper()[8:12]}-{uid.hex.upper()[12:16]}-{uid.hex.upper()[16:20]}-{uid.hex.upper()[20:32]}"
    
    def _get_mime_type(self, source_path: Path) -> str:
        """Get MIME type from file extension"""
        mime_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.md': 'text/markdown',
            '.txt': 'text/plain',
            '.json': 'application/json',
            '.xml': 'application/xml'
        }
        return mime_map.get(source_path.suffix.lower(), 'application/octet-stream')
    
    def _get_extraction_method(self, source_path: Path) -> str:
        """Get extraction method name"""
        method_map = {
            '.pdf': 'pdfplumber',
            '.docx': 'python-docx',
            '.xlsx': 'openpyxl',
            '.pptx': 'python-pptx',
            '.md': 'builtin',
            '.txt': 'builtin',
            '.json': 'builtin',
            '.xml': 'builtin'
        }
        return method_map.get(source_path.suffix.lower(), 'unknown')
    
    def _persist_evidence_pack(self, evidence_pack: Dict[str, Any]) -> str:
        """
        Persist Evidence Pack to PostgreSQL (append-only)
        
        Returns reference to persisted pack
        """
        if not self.postgres_agent:
            # Fallback: return mock reference
            return f"mock://evidence_packs/{evidence_pack['evidence_id']}"
        
        try:
            with self.postgres_agent.connection.cursor() as cur:
                # Build column list dynamically — tenant_id is optional
                columns = [
                    "evidence_id", "chunk_id", "schema_version", "created_utc",
                    "source_ref", "normalized_text", "technical_metadata",
                    "integrity", "sampling_policy_ref", "tags",
                ]
                values = [
                    evidence_pack['evidence_id'],
                    evidence_pack['chunk_id'],
                    evidence_pack['schema_version'],
                    evidence_pack['created_utc'],
                    json.dumps(evidence_pack['source_ref']),
                    evidence_pack['normalized_text'],
                    json.dumps(evidence_pack['technical_metadata']),
                    json.dumps(evidence_pack['integrity']),
                    evidence_pack.get('sampling_policy_ref'),
                    json.dumps(evidence_pack.get('tags', [])),
                ]
                if evidence_pack.get('tenant_id'):
                    columns.append("tenant_id")
                    values.append(evidence_pack['tenant_id'])
                
                placeholders = ", ".join(["%s"] * len(values))
                col_names = ", ".join(columns)
                cur.execute(
                    f"INSERT INTO evidence_packs ({col_names}) VALUES ({placeholders}) RETURNING id",
                    values,
                )
                row_id = cur.fetchone()[0]
            self.postgres_agent.connection.commit()
            return f"postgres://evidence_packs/{row_id}"
        except Exception as e:
            logger.error(f"Failed to persist Evidence Pack: {e}")
            return f"error://failed_to_persist/{evidence_pack['evidence_id']}"
